"""Local file loader - recursively parses files from a directory."""

import json
from dataclasses import dataclass, field
from pathlib import Path

from rich.console import Console

from kbskills.utils.text import clean_text

console = Console()


@dataclass
class Document:
    source: str
    content: str
    metadata: dict = field(default_factory=dict)


# Extensions that we can handle
TEXT_EXTENSIONS = {".md", ".txt", ".json", ".jsonl", ".csv", ".tsv", ".yaml", ".yml", ".xml", ".html", ".htm"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}


def load_directory(dir_path: str | Path) -> list[Document]:
    """Recursively load all supported files from a directory."""
    dir_path = Path(dir_path)
    if not dir_path.is_dir():
        raise FileNotFoundError(f"Directory not found: {dir_path}")

    documents = []
    for file_path in sorted(dir_path.rglob("*")):
        if file_path.is_dir():
            continue
        doc = load_file(file_path)
        if doc:
            documents.append(doc)

    console.print(f"[green]Loaded {len(documents)} documents from {dir_path}[/green]")
    return documents


def load_file(file_path: Path) -> Document | None:
    """Load a single file and return a Document, or None if unsupported."""
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _load_pdf(file_path)
        elif ext == ".docx":
            return _load_docx(file_path)
        elif ext == ".pptx":
            return _load_pptx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _load_excel(file_path)
        elif ext == ".csv":
            return _load_csv(file_path)
        elif ext in TEXT_EXTENSIONS:
            return _load_text(file_path)
        elif ext in IMAGE_EXTENSIONS:
            # Image files require Gemini Vision - skip for now, handle in pipeline
            return Document(
                source=str(file_path),
                content="",
                metadata={"type": "image", "needs_vision": True},
            )
        else:
            return None
    except Exception as e:
        console.print(f"[yellow]Warning: Failed to load {file_path}: {e}[/yellow]")
        return None


def _load_text(file_path: Path) -> Document:
    content = file_path.read_text(encoding="utf-8", errors="replace")
    return Document(
        source=str(file_path),
        content=clean_text(content),
        metadata={"type": file_path.suffix.lstrip(".")},
    )


def _load_pdf(file_path: Path) -> Document:
    import pymupdf

    doc = pymupdf.open(str(file_path))
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()

    return Document(
        source=str(file_path),
        content=clean_text("\n\n".join(pages)),
        metadata={"type": "pdf", "pages": len(pages)},
    )


def _load_docx(file_path: Path) -> Document:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    return Document(
        source=str(file_path),
        content=clean_text("\n\n".join(paragraphs)),
        metadata={"type": "docx"},
    )


def _load_pptx(file_path: Path) -> Document:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    return Document(
        source=str(file_path),
        content=clean_text("\n\n".join(texts)),
        metadata={"type": "pptx", "slides": len(prs.slides)},
    )


def _load_excel(file_path: Path) -> Document:
    import pandas as pd

    sheets = pd.read_excel(str(file_path), sheet_name=None)
    parts = []
    for name, df in sheets.items():
        parts.append(f"[Sheet: {name}]\n{df.to_string()}")

    return Document(
        source=str(file_path),
        content=clean_text("\n\n".join(parts)),
        metadata={"type": "excel", "sheets": list(sheets.keys())},
    )


def _load_csv(file_path: Path) -> Document:
    import pandas as pd

    df = pd.read_csv(str(file_path))
    return Document(
        source=str(file_path),
        content=clean_text(df.to_string()),
        metadata={"type": "csv", "rows": len(df)},
    )
